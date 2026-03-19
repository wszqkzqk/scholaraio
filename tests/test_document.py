"""Tests for scholaraio.document (Office document inspection)."""

from __future__ import annotations

from pathlib import Path

import pytest

pytest.importorskip("docx")
pytest.importorskip("pptx")
pytest.importorskip("openpyxl")


# ---------------------------------------------------------------------------
# Fixtures — create minimal Office files in tmp_path
# ---------------------------------------------------------------------------


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.paragraphs[0].text = "Hello World"

    # Add a table
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide2.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(5), Inches(3)).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def sample_docx(tmp_path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a test paragraph.")
    table = doc.add_table(rows=2, cols=3)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    table.rows[0].cells[2].text = "C"
    doc.add_heading("Section 2", level=2)
    doc.add_paragraph("Another paragraph.")

    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


@pytest.fixture()
def sample_xlsx(tmp_path: Path) -> Path:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Score", "Grade"])
    ws.append(["Alice", 95, "A"])
    ws.append(["Bob", 87, "B"])
    ws.freeze_panes = "A2"

    chart = BarChart()
    chart.title = "Scores"
    data = Reference(ws, min_col=2, min_row=1, max_row=3)
    chart.add_data(data, titles_from_data=True)
    ws.add_chart(chart, "E2")

    ws2 = wb.create_sheet("Summary")
    ws2.append(["Total", 182])

    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestInspectDispatcher:
    def test_auto_detect_pptx(self, sample_pptx):
        from scholaraio.document import inspect

        result = inspect(sample_pptx)
        assert "PPTX" in result

    def test_auto_detect_docx(self, sample_docx):
        from scholaraio.document import inspect

        result = inspect(sample_docx)
        assert "DOCX" in result

    def test_auto_detect_xlsx(self, sample_xlsx):
        from scholaraio.document import inspect

        result = inspect(sample_xlsx)
        assert "XLSX" in result

    def test_format_override(self, sample_pptx):
        from scholaraio.document import inspect

        result = inspect(sample_pptx, fmt="pptx")
        assert "PPTX" in result

    def test_unsupported_format(self, tmp_path):
        from scholaraio.document import inspect

        p = tmp_path / "test.pdf"
        p.write_text("dummy")
        with pytest.raises(ValueError, match="不支持"):
            inspect(p)

    def test_file_not_found(self, tmp_path):
        from scholaraio.document import inspect

        with pytest.raises(FileNotFoundError):
            inspect(tmp_path / "nonexistent.pptx")

    def test_directory_path_rejected(self, tmp_path):
        from scholaraio.document import inspect

        with pytest.raises(ValueError, match="不是文件"):
            inspect(tmp_path)


class TestInspectPptx:
    def test_slide_count(self, sample_pptx):
        from scholaraio.document import inspect_pptx

        result = inspect_pptx(sample_pptx)
        assert "2 页" in result

    def test_text_content(self, sample_pptx):
        from scholaraio.document import inspect_pptx

        result = inspect_pptx(sample_pptx)
        assert "Hello World" in result

    def test_table_detected(self, sample_pptx):
        from scholaraio.document import inspect_pptx

        result = inspect_pptx(sample_pptx)
        assert "[Table 3x2]" in result
        assert "Name" in result

    def test_summary_section(self, sample_pptx):
        from scholaraio.document import inspect_pptx

        result = inspect_pptx(sample_pptx)
        assert "总结" in result
        assert "表格: 1" in result


class TestInspectDocx:
    def test_heading_detected(self, sample_docx):
        from scholaraio.document import inspect_docx

        result = inspect_docx(sample_docx)
        assert "Test Document" in result
        assert "Heading 1" in result

    def test_paragraph_detected(self, sample_docx):
        from scholaraio.document import inspect_docx

        result = inspect_docx(sample_docx)
        assert "test paragraph" in result

    def test_table_detected(self, sample_docx):
        from scholaraio.document import inspect_docx

        result = inspect_docx(sample_docx)
        assert "[Table 2x3]" in result

    def test_summary(self, sample_docx):
        from scholaraio.document import inspect_docx

        result = inspect_docx(sample_docx)
        assert "总结" in result
        assert "H1:" in result
        assert "H2:" in result


class TestInspectXlsx:
    def test_sheet_names(self, sample_xlsx):
        from scholaraio.document import inspect_xlsx

        result = inspect_xlsx(sample_xlsx)
        assert "Data" in result
        assert "Summary" in result

    def test_header_preview(self, sample_xlsx):
        from scholaraio.document import inspect_xlsx

        result = inspect_xlsx(sample_xlsx)
        assert "Name" in result
        assert "Score" in result

    def test_data_preview(self, sample_xlsx):
        from scholaraio.document import inspect_xlsx

        result = inspect_xlsx(sample_xlsx)
        assert "Alice" in result

    def test_frozen_panes(self, sample_xlsx):
        from scholaraio.document import inspect_xlsx

        result = inspect_xlsx(sample_xlsx)
        assert "冻结窗格" in result

    def test_chart_detected(self, sample_xlsx):
        from scholaraio.document import inspect_xlsx

        result = inspect_xlsx(sample_xlsx)
        assert "BarChart" in result
        assert "Scores" in result

    def test_summary(self, sample_xlsx):
        from scholaraio.document import inspect_xlsx

        result = inspect_xlsx(sample_xlsx)
        assert "图表: 1" in result
