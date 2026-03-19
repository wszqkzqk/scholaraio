"""Office 文档检查工具（DOCX / PPTX / XLSX）。

Provides structured text summaries of Office documents for AI agents
to verify layout, content, and detect common issues (overflow, missing
content, etc.) without opening a GUI application.
"""

from __future__ import annotations

import logging
from pathlib import Path

_log = logging.getLogger(__name__)

# EMU per inch constant (used by python-pptx)
_EMU_PER_INCH = 914400


# ---------------------------------------------------------------------------
# Public dispatcher
# ---------------------------------------------------------------------------


def inspect(path: Path, fmt: str | None = None) -> str:
    """Inspect an Office document and return a structured text summary.

    Args:
        path: Path to the Office file.
        fmt: File format override (``"docx"``, ``"pptx"``, ``"xlsx"``).
            Auto-detected from extension when *None*.

    Returns:
        Human-readable inspection report.

    Raises:
        ValueError: If the file format is unsupported or cannot be detected.
        FileNotFoundError: If *path* does not exist.
    """
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")
    if not path.is_file():
        raise ValueError(f"路径不是文件: {path}")

    if fmt is None:
        fmt = path.suffix.lstrip(".").lower()

    dispatch = {
        "pptx": inspect_pptx,
        "docx": inspect_docx,
        "xlsx": inspect_xlsx,
    }
    func = dispatch.get(fmt)
    if func is None:
        raise ValueError(f"不支持的文件格式: .{fmt}（支持 docx / pptx / xlsx）")
    return func(path)


# ---------------------------------------------------------------------------
# PPTX inspection
# ---------------------------------------------------------------------------


def inspect_pptx(path: Path) -> str:
    """Inspect a PowerPoint file and return a structured text summary.

    Args:
        path: Path to the ``.pptx`` file.

    Returns:
        Slide-by-slide report with shape positions, sizes, text preview,
        image info, and overflow warnings.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")

    prs = Presentation(str(path))
    sw = (prs.slide_width or 0) / _EMU_PER_INCH
    sh = (prs.slide_height or 0) / _EMU_PER_INCH
    total = len(prs.slides)

    lines: list[str] = []
    lines.append(f"=== PPTX: {path.name} ===")
    lines.append(f'幻灯片尺寸: {sw:.1f}" x {sh:.1f}" ({total} 页)')
    lines.append("")

    for i, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name
        lines.append(f"--- Slide {i}/{total} (layout: {layout_name}) ---")

        warnings: list[str] = []

        for shape in slide.shapes:
            left = shape.left / _EMU_PER_INCH if shape.left else 0
            top = shape.top / _EMU_PER_INCH if shape.top else 0
            w = shape.width / _EMU_PER_INCH if shape.width else 0
            h = shape.height / _EMU_PER_INCH if shape.height else 0
            right = left + w
            bottom = top + h

            # Overflow detection
            if right > sw + 0.1:
                warnings.append(f'  \u26a0 溢出: shape 右边界 {right:.1f}" > 幻灯片宽度 {sw:.1f}"')
            if bottom > sh + 0.1:
                warnings.append(f'  \u26a0 溢出: shape 下边界 {bottom:.1f}" > 幻灯片高度 {sh:.1f}"')

            pos = f'({left:.1f}",{top:.1f}") {w:.1f}"x{h:.1f}"'

            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    size_kb = len(img.blob) // 1024
                    lines.append(f"  [Image] {pos}  [{img.content_type}, {size_kb}KB]")
                except Exception:
                    lines.append(f"  [Image] {pos}")
                continue

            # Table
            if shape.has_table:
                tbl = shape.table
                nr = len(tbl.rows)
                nc = len(tbl.columns)
                header_cells = [tbl.cell(0, c).text[:20] for c in range(min(nc, 5))]
                hdr_str = " | ".join(header_cells)
                if nc > 5:
                    hdr_str += " | ..."
                lines.append(f"  [Table {nr}x{nc}] {pos}")
                lines.append(f"    表头: {hdr_str}")
                continue

            # Text frame
            if shape.has_text_frame:
                tf = shape.text_frame
                paras = [p.text.strip() for p in tf.paragraphs if p.text.strip()]
                if not paras:
                    continue

                # Determine shape label
                label = "TextBox"
                try:
                    if shape.placeholder_format is not None:
                        idx = shape.placeholder_format.idx
                        label = f"Placeholder {idx}"
                except ValueError:
                    pass  # not a placeholder

                lines.append(f"  [{label}] {pos}")

                # Show first 3 paragraphs, truncated
                for p_text in paras[:3]:
                    display = p_text[:60] + ("..." if len(p_text) > 60 else "")
                    lines.append(f'    "{display}"')
                if len(paras) > 3:
                    lines.append(f"    ... (+{len(paras) - 3} more)")

                # Font info from first paragraph
                first_para = tf.paragraphs[0]
                if first_para.runs:
                    run = first_para.runs[0]
                    font_info = []
                    if run.font.size:
                        font_info.append(f"{run.font.size.pt:.0f}pt")
                    if run.font.bold:
                        font_info.append("bold")
                    if run.font.italic:
                        font_info.append("italic")
                    if font_info:
                        lines.append(f"    font: {', '.join(font_info)}")

                # Estimate text overflow (rough heuristic)
                # Accumulate per-paragraph height to avoid bias from mixed font sizes.
                est_height = 0.0
                for para in tf.paragraphs:
                    font_size = 18  # default pt
                    if para.runs and para.runs[0].font.size:
                        font_size = para.runs[0].font.size.pt
                    line_height_in = 1.2 * (font_size / 72) if font_size else 0.3
                    text_len = len(para.text)
                    if text_len == 0:
                        est_height += line_height_in * 0.5  # empty para still takes space
                        continue
                    # Estimate chars per line based on shape width and font size
                    chars_per_line = max(1, int(w * 72 / (font_size * 0.55)))
                    est_lines = max(1, -(-text_len // chars_per_line))  # ceil div
                    est_height += est_lines * line_height_in
                if est_height > h * 1.1 and h > 0.5:
                    warnings.append(f'  \u26a0 文字可能溢出: 估算高度 {est_height:.1f}" > 容器高度 {h:.1f}"')

        if warnings:
            for w_line in warnings:
                lines.append(w_line)

        lines.append("")

    # Summary
    img_count = sum(1 for slide in prs.slides for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    tbl_count = sum(1 for slide in prs.slides for s in slide.shapes if s.has_table)
    lines.append("--- 总结 ---")
    lines.append(f"页数: {total} | 图片: {img_count} | 表格: {tbl_count}")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# DOCX inspection
# ---------------------------------------------------------------------------


def inspect_docx(path: Path) -> str:
    """Inspect a Word document and return a structured text summary.

    Args:
        path: Path to the ``.docx`` file.

    Returns:
        Document structure report with headings, paragraphs, tables,
        images, and style info.
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        raise ImportError("python-docx 未安装，请运行: pip install python-docx")

    doc = Document(str(path))

    lines: list[str] = []
    lines.append(f"=== DOCX: {path.name} ===")

    # Section info
    for si, section in enumerate(doc.sections):
        pw = section.page_width.inches if section.page_width else 0
        ph = section.page_height.inches if section.page_height else 0
        orient = "横向" if pw > ph else "纵向"
        lines.append(f'  Section {si + 1}: {pw:.1f}" x {ph:.1f}" ({orient})')

    # Count elements and collect styles
    styles_used: set[str] = set()
    h_counts: dict[str, int] = {}
    para_count = 0
    table_count = 0
    image_count = 0

    lines.append("")
    lines.append("--- 文档结构 ---")

    # Walk document body in order (block-level) to avoid paragraph/table desync.
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)

            style_name = para.style.name if para.style else "Normal"
            styles_used.add(style_name)
            text = para.text.strip()

            # Check for images in this paragraph
            drawings = child.findall(f".//{qn('w:drawing')}")
            if drawings:
                image_count += len(drawings)
                for _ in drawings:
                    lines.append("  [Image] (嵌入图片)")

            if style_name.startswith("Heading"):
                level = style_name.replace("Heading ", "H")
                h_counts[level] = h_counts.get(level, 0) + 1
                display = text[:70] + ("..." if len(text) > 70 else "")
                lines.append(f'  [{style_name}] "{display}"')
            elif style_name.startswith("List"):
                # Emit each list paragraph as one list-item line.
                lines.append(f'  [{style_name}] "{text[:60]}"')
                para_count += 1
            elif "toc" in style_name.lower():
                if para_count == 0 or not lines[-1].startswith("  [TOC"):
                    lines.append("  [TOC] (目录字段)")
                para_count += 1
            elif text:
                para_count += 1
                display = text[:70] + ("..." if len(text) > 70 else "")
                # Font info
                font_info = ""
                if para.runs:
                    run = para.runs[0]
                    parts = []
                    if run.font.name:
                        parts.append(run.font.name)
                    if run.font.size:
                        parts.append(f"{run.font.size.pt:.0f}pt")
                    if run.bold:
                        parts.append("bold")
                    if parts:
                        font_info = f" [{', '.join(parts)}]"
                lines.append(f'  [Paragraph] "{display}"{font_info}')
            else:
                # Empty paragraph (spacer)
                para_count += 1

        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            table_count += 1

            nr = len(table.rows)
            nc = len(table.columns)
            style = table.style.name if table.style else "?"
            lines.append(f"  [Table {nr}x{nc}] (style: {style})")

            # Header row preview
            if nr > 0:
                headers = [table.rows[0].cells[c].text[:20] for c in range(min(nc, 5))]
                hdr_str = " | ".join(headers)
                if nc > 5:
                    hdr_str += " | ..."
                lines.append(f"    表头: {hdr_str}")
                lines.append(f"    ({nr - 1} 数据行)")

        elif child.tag == qn("w:sectPr"):
            pass  # section properties, already handled

    # Summary
    lines.append("")
    lines.append("--- 总结 ---")
    h_str = ", ".join(f"{k}:{v}" for k, v in sorted(h_counts.items()))
    lines.append(
        f"段落: {para_count} | 表格: {table_count} | 图片: {image_count} | 标题: {sum(h_counts.values())} ({h_str})"
    )
    lines.append(f"样式: {', '.join(sorted(styles_used))}")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# XLSX inspection
# ---------------------------------------------------------------------------


def inspect_xlsx(path: Path) -> str:
    """Inspect an Excel workbook and return a structured text summary.

    Args:
        path: Path to the ``.xlsx`` file.

    Returns:
        Sheet-by-sheet overview with dimensions, data preview, charts,
        and formatting info.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl 未安装，请运行: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), read_only=False, data_only=True)
    try:
        lines: list[str] = []
        lines.append(f"=== XLSX: {path.name} ===")
        lines.append(f"工作表: {len(wb.sheetnames)} ({', '.join(wb.sheetnames)})")
        lines.append("")

        total_rows = 0

        for ws_name in wb.sheetnames:
            ws = wb[ws_name]
            active_mark = " (active)" if ws == wb.active else ""
            lines.append(f'--- Sheet "{ws_name}"{active_mark} ---')

            dims = ws.dimensions or "?"
            mr = ws.max_row or 0
            mc = ws.max_column or 0
            total_rows += mr
            lines.append(f"  范围: {dims} ({mr} 行 x {mc} 列)")

            # Frozen panes
            if ws.freeze_panes:
                lines.append(f"  冻结窗格: {ws.freeze_panes}")

            # Auto-filter
            if ws.auto_filter and ws.auto_filter.ref:
                lines.append(f"  自动筛选: {ws.auto_filter.ref}")

            # Merged cells
            if ws.merged_cells.ranges:
                merges = [str(r) for r in list(ws.merged_cells.ranges)[:5]]
                merge_str = ", ".join(merges)
                if len(ws.merged_cells.ranges) > 5:
                    merge_str += f" ... (+{len(ws.merged_cells.ranges) - 5})"
                lines.append(f"  合并单元格: {merge_str}")

            # Header row (row 1)
            if mr > 0 and mc > 0:
                headers = []
                for c in range(1, min(mc + 1, 8)):
                    val = ws.cell(row=1, column=c).value
                    headers.append(str(val)[:25] if val is not None else "")
                hdr_str = " | ".join(headers)
                if mc > 7:
                    hdr_str += " | ..."
                lines.append(f"  表头 (row 1): {hdr_str}")

            # Data preview (rows 2-4)
            preview_rows = min(mr, 4) - 1
            if preview_rows > 0:
                lines.append("  数据预览:")
                for r in range(2, 2 + preview_rows):
                    vals = []
                    for c in range(1, min(mc + 1, 8)):
                        val = ws.cell(row=r, column=c).value
                        s = str(val)[:25] if val is not None else ""
                        vals.append(s)
                    lines.append(f"    Row {r}: {' | '.join(vals)}")

            # Charts
            if ws._charts:
                lines.append(f"  图表: {len(ws._charts)}")
                for ci, chart in enumerate(ws._charts, 1):
                    chart_type = type(chart).__name__
                    # Extract title string from openpyxl chart title object
                    title_str = "(无标题)"
                    if chart.title:
                        if isinstance(chart.title, str):
                            title_str = chart.title
                        else:
                            # openpyxl Title object — extract text from rich text runs
                            try:
                                for para in chart.title.tx.rich.paragraphs:
                                    for run in para.r:
                                        if run.t:
                                            title_str = run.t
                                            break
                                    if title_str != "(无标题)":
                                        break
                            except (AttributeError, TypeError):
                                title_str = str(chart.title)[:40]
                    lines.append(f'    [{ci}] {chart_type}: "{title_str}"')

            lines.append("")

        # Summary
        chart_total = sum(len(wb[s]._charts) for s in wb.sheetnames)
        lines.append("--- 总结 ---")
        lines.append(f"工作表: {len(wb.sheetnames)} | 总行数: ~{total_rows} | 图表: {chart_total}")

        return "\n".join(lines)
    finally:
        wb.close()
